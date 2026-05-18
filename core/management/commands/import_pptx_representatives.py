"""PPTX faylidan vakillarni import qiladi.

Foydalanish:
    python manage.py import_pptx_representatives <PPTX_PATH> <direction_key>
    python manage.py import_pptx_representatives /tmp/teatr.pptx theater_circus --clear

Opsiyalar:
    --clear      Avval shu yo'nalishdagi vakillarni o'chiradi
    --dry-run    DB ga yozmasdan tahlil qiladi
    --limit N    Birinchi N ta slayd
"""
import re
import sys
from datetime import datetime
from pathlib import Path

from django.core.files.base import ContentFile
from django.core.management.base import BaseCommand
from django.db import transaction

from core.models import (
    AwardName,
    Direction,
    FamilyMember,
    Representative,
    RepresentativeAward,
)


# ─── Cyrillic → Latin (O'zbek standarti) ──────────────────────────────────

CYR_TO_LATN = [
    ('Ё', 'Yo'), ('ё', 'yo'),
    ('Ю', 'Yu'), ('ю', 'yu'),
    ('Я', 'Ya'), ('я', 'ya'),
    ('Ч', 'Ch'), ('ч', 'ch'),
    ('Ш', 'Sh'), ('ш', 'sh'),
    ('Ў', "O'"), ('ў', "o'"),
    ('Қ', 'Q'), ('қ', 'q'),
    ('Ғ', "G'"), ('ғ', "g'"),
    ('Ҳ', 'H'), ('ҳ', 'h'),
    ('Ж', 'J'), ('ж', 'j'),
    ('Ц', 'Ts'), ('ц', 'ts'),
    ('Й', 'Y'), ('й', 'y'),
    ('Э', 'E'), ('э', 'e'),
    ('А', 'A'), ('а', 'a'),
    ('Б', 'B'), ('б', 'b'),
    ('В', 'V'), ('в', 'v'),
    ('Г', 'G'), ('г', 'g'),
    ('Д', 'D'), ('д', 'd'),
    ('Е', 'E'), ('е', 'e'),
    ('З', 'Z'), ('з', 'z'),
    ('И', 'I'), ('и', 'i'),
    ('К', 'K'), ('к', 'k'),
    ('Л', 'L'), ('л', 'l'),
    ('М', 'M'), ('м', 'm'),
    ('Н', 'N'), ('н', 'n'),
    ('О', 'O'), ('о', 'o'),
    ('П', 'P'), ('п', 'p'),
    ('Р', 'R'), ('р', 'r'),
    ('С', 'S'), ('с', 's'),
    ('Т', 'T'), ('т', 't'),
    ('У', 'U'), ('у', 'u'),
    ('Ф', 'F'), ('ф', 'f'),
    ('Х', 'X'), ('х', 'x'),
    ('Ъ', "'"), ('ъ', "'"),
    ('Ь', ''), ('ь', ''),
    ('Ы', 'I'), ('ы', 'i'),
    ('Щ', 'Shch'), ('щ', 'shch'),
]


def to_latn(text):
    if not text:
        return ''
    for cyr, lat in CYR_TO_LATN:
        text = text.replace(cyr, lat)
    return text


def clean(text):
    if not text:
        return ''
    text = text.replace('\x0b', ' ').replace(' ', ' ').strip()
    text = re.sub(r' +', ' ', text)
    return text


def parse_date(text):
    if not text:
        return None
    text = clean(text)
    for fmt in ('%d.%m.%Y', '%d/%m/%Y', '%Y-%m-%d'):
        try:
            return datetime.strptime(text, fmt).date()
        except ValueError:
            continue
    return None


NATIONALITY_MAP = {
    'ўзбек': 'ozbek', 'узбек': 'ozbek', "o'zbek": 'ozbek',
    'қозоқ': 'qozoq', 'казах': 'qozoq',
    'тожик': 'tojik', 'таджик': 'tojik',
    'рус': 'rus',
    'қорақалпоқ': 'qoraqalpoq', 'каракалпак': 'qoraqalpoq',
    'қирғиз': 'qirgiz', 'киргиз': 'qirgiz',
    'туркман': 'turkman',
    'татар': 'tatar',
    'корейс': 'koreys', 'кореец': 'koreys',
}


def map_nationality(text):
    if not text:
        return ''
    return NATIONALITY_MAP.get(clean(text).lower(), '')


RELATION_MAP = {
    'отаси': FamilyMember.RELATION_FATHER, 'отец': FamilyMember.RELATION_FATHER,
    'онаси': FamilyMember.RELATION_MOTHER, 'мать': FamilyMember.RELATION_MOTHER,
    'турмуш ўртоғи': FamilyMember.RELATION_SPOUSE,
    'хотини': FamilyMember.RELATION_SPOUSE, 'жена': FamilyMember.RELATION_SPOUSE,
    'эри': FamilyMember.RELATION_SPOUSE, 'муж': FamilyMember.RELATION_SPOUSE,
    'фарзанди': FamilyMember.RELATION_CHILD, 'ўғли': FamilyMember.RELATION_CHILD,
    'қизи': FamilyMember.RELATION_CHILD,
    'сын': FamilyMember.RELATION_CHILD, 'дочь': FamilyMember.RELATION_CHILD,
}


def map_relation(text):
    if not text:
        return FamilyMember.RELATION_OTHER
    key = clean(text).rstrip(':').lower()
    for k, v in RELATION_MAP.items():
        if k in key:
            return v
    return FamilyMember.RELATION_OTHER


def cell_text(table, r, c):
    try:
        return clean(table.cell(r, c).text_frame.text)
    except Exception:
        return ''


def extract_picture(slide):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            try:
                img = shape.image
                return img.blob, img.ext
            except Exception:
                pass
    return None, None


def parse_slide(slide, idx):
    tables = [s.table for s in slide.shapes if s.has_table]
    if len(tables) < 2:
        return None

    data = {
        'slide_idx': idx,
        'last_name': '', 'first_name': '', 'middle_name': '',
        'nationality': '', 'birth_date': None, 'birth_place': '',
        'marital_status': '', 'family': [],
        'university': '', 'specialty': '', 'academic_degree': '',
        'foreign_lang': '', 'training': '',
        'position': '', 'career_level': '', 'total_experience': '',
        'leadership_experience': '', 'leadership_positions': '',
        'health_text': '', 'last_medical_treatment': '',
        'medical_checkup': '', 'health_problems': '',
        'awards_raw': '',
        'description': '', 'state_events': '',
    }

    # Table 1: Shaxsiy
    t1 = tables[0]
    data['last_name'] = cell_text(t1, 1, 4)
    data['first_name'] = cell_text(t1, 2, 4)
    data['middle_name'] = cell_text(t1, 3, 4)
    data['nationality'] = cell_text(t1, 4, 4)
    data['birth_date'] = parse_date(cell_text(t1, 5, 4))
    data['birth_place'] = cell_text(t1, 6, 4)
    data['marital_status'] = cell_text(t1, 8, 2)

    # Oila a'zolari
    for ri in range(10, 13):
        rel_label = cell_text(t1, ri, 0)
        name_info = cell_text(t1, ri, 1)
        note = cell_text(t1, ri, 4)
        if not name_info and not note:
            continue
        parts = re.split(r'[\n,]', name_info, maxsplit=1)
        name = clean(parts[0]) if parts else ''
        info = clean(parts[1]) if len(parts) > 1 else ''
        data['family'].append({
            'relation': map_relation(rel_label),
            'name': name,
            'info': info,
            'note': clean(note),
        })

    # Ta'lim (col 3 da)
    data['university'] = cell_text(t1, 14, 3) or cell_text(t1, 14, 2)
    data['specialty'] = cell_text(t1, 15, 3) or cell_text(t1, 15, 2)
    data['academic_degree'] = cell_text(t1, 16, 3) or cell_text(t1, 16, 2)
    data['foreign_lang'] = cell_text(t1, 17, 3) or cell_text(t1, 17, 2)
    data['training'] = cell_text(t1, 18, 3) or cell_text(t1, 18, 2)

    # Table 2: Mehnat faoliyati
    t2 = tables[1]
    data['position'] = cell_text(t2, 1, 1)
    data['career_level'] = cell_text(t2, 2, 1)
    data['total_experience'] = cell_text(t2, 3, 1)
    data['leadership_experience'] = cell_text(t2, 4, 1)
    data['leadership_positions'] = cell_text(t2, 5, 1)
    data['health_text'] = cell_text(t2, 6, 1)
    data['last_medical_treatment'] = cell_text(t2, 7, 1)
    data['medical_checkup'] = cell_text(t2, 8, 1)
    data['health_problems'] = cell_text(t2, 9, 1)

    # Mukofotlar
    for ri in range(len(t2.rows) - 1, 9, -1):
        txt = cell_text(t2, ri, 0)
        if 'мукофот' in txt.lower() or 'фахрий' in txt.lower() or ' й.' in txt or ' й ' in txt:
            data['awards_raw'] = cell_text(t2, ri, 0)
            break

    # Table 3: Faoliyat
    if len(tables) >= 3:
        t3 = tables[2]
        for ri in range(len(t3.rows)):
            txt = cell_text(t3, ri, 0)
            if not txt:
                continue
            low = txt.lower()
            if 'қисқача тавсифнома' in low or 'қисқача тавсиф' in low:
                if ri + 1 < len(t3.rows):
                    data['description'] = cell_text(t3, ri + 1, 0)
            elif 'давлат тадбирлари' in low:
                if ri + 1 < len(t3.rows):
                    data['state_events'] = cell_text(t3, ri + 1, 0)

    # Health: text → choice
    data['health'] = ''
    htxt = data['health_text'].lower()
    if 'яхши' in htxt:
        data['health'] = Representative.HEALTH_GOOD
    elif 'ўрта' in htxt or 'qoniqarli' in htxt:
        data['health'] = Representative.HEALTH_AVERAGE
    elif 'ёмон' in htxt or 'қониқарсиз' in htxt:
        data['health'] = Representative.HEALTH_POOR

    return data


AWARD_LINE_RE = re.compile(r"(\d{4})\s*(?:й\.?|йил|год|г\.?)?\s*[\.\-—–]?\s*(.+?)$", re.IGNORECASE)


def parse_awards(text, award_lookup):
    items = []
    if not text:
        return items
    for line in text.split('\n'):
        line = clean(line)
        if not line:
            continue
        m = AWARD_LINE_RE.match(line)
        if not m:
            continue
        year = int(m.group(1))
        name_text = clean(m.group(2)).strip('-—–.').strip()
        if not name_text:
            continue
        matched = None
        nm_lower = name_text.lower()
        for key, award_obj in award_lookup.items():
            if key in nm_lower or nm_lower in key:
                matched = award_obj
                break
        items.append({'year': year, 'award': matched, 'raw_name': name_text})
    return items


class Command(BaseCommand):
    help = "PPTX faylidan vakillarni import qiladi"

    def add_arguments(self, parser):
        parser.add_argument('pptx_path', type=str)
        parser.add_argument('direction_key', type=str)
        parser.add_argument('--clear', action='store_true')
        parser.add_argument('--dry-run', action='store_true')
        parser.add_argument('--limit', type=int, default=0)

    def handle(self, *args, **options):
        if hasattr(sys.stdout, 'reconfigure'):
            try:
                sys.stdout.reconfigure(encoding='utf-8')
            except Exception:
                pass

        from pptx import Presentation

        pptx_path = Path(options['pptx_path']).resolve()
        if not pptx_path.exists():
            self.stderr.write(f'Fayl topilmadi: {pptx_path}')
            return

        try:
            direction = Direction.objects.get(key=options['direction_key'])
        except Direction.DoesNotExist:
            self.stderr.write(f"Yo'nalish topilmadi: {options['direction_key']}")
            return

        if options['clear'] and not options['dry_run']:
            n = Representative.objects.filter(direction=direction).count()
            Representative.objects.filter(direction=direction).delete()
            self.stdout.write(self.style.WARNING(f"O'chirildi: {n} ta vakil"))

        award_lookup = {}
        for a in AwardName.objects.all():
            key = (a.name_uz_cyrl or a.name).lower()
            key = re.sub(r'[«»“”"\'.,]', '', key).strip()
            award_lookup[key] = a

        pres = Presentation(str(pptx_path))
        total = len(pres.slides)
        limit = options['limit'] or total
        self.stdout.write(f'Slaydlar: {total}, import qilinadi: {min(limit, total)}')

        created = skipped = failed = 0
        unmatched_awards = []

        for idx, slide in enumerate(pres.slides):
            if idx >= limit:
                break
            try:
                data = parse_slide(slide, idx)
                if not data or not data['last_name']:
                    skipped += 1
                    continue

                if options['dry_run']:
                    self.stdout.write(
                        f"  [{idx+1}/{total}] DRY {data['last_name']} {data['first_name']} {data['middle_name']}"
                    )
                    continue

                with transaction.atomic():
                    rep = self._create_representative(direction, data)
                    awards = parse_awards(data['awards_raw'], award_lookup)
                    for a in awards:
                        if a['award']:
                            RepresentativeAward.objects.create(
                                representative=rep, award=a['award'], year=a['year']
                            )
                        else:
                            unmatched_awards.append((rep.full_name, a['year'], a['raw_name']))

                    img_blob, img_ext = extract_picture(slide)
                    if img_blob:
                        filename = f'pptx_{idx+1}.{img_ext or "jpg"}'
                        rep.photo.save(filename, ContentFile(img_blob), save=True)

                created += 1
                self.stdout.write(f"  [{idx+1}/{total}] +{rep.full_name}")

            except Exception as e:
                failed += 1
                self.stderr.write(f"  [{idx+1}/{total}] FAIL: {e}")

        self.stdout.write(self.style.SUCCESS(
            f"\n{created} qo'shildi, {skipped} o'tkazib yuborilgan, {failed} xato."
        ))

        if unmatched_awards:
            self.stdout.write(self.style.WARNING(
                f"\nMos kelmagan mukofotlar: {len(unmatched_awards)} ta"
            ))

    def _create_representative(self, direction, d):
        ln, fn, mn = d['last_name'], d['first_name'], d['middle_name']
        rep = Representative.objects.create(
            direction=direction,
            last_name=to_latn(ln), last_name_uz_latn=to_latn(ln),
            last_name_uz_cyrl=ln, last_name_ru=ln,
            first_name=to_latn(fn), first_name_uz_latn=to_latn(fn),
            first_name_uz_cyrl=fn, first_name_ru=fn,
            middle_name=to_latn(mn), middle_name_uz_latn=to_latn(mn),
            middle_name_uz_cyrl=mn, middle_name_ru=mn,
            gender=Representative.GENDER_MALE,
            nationality=map_nationality(d['nationality']),
            birth_date=d['birth_date'],
            birth_place=to_latn(d['birth_place']), birth_place_uz_latn=to_latn(d['birth_place']),
            birth_place_uz_cyrl=d['birth_place'], birth_place_ru=d['birth_place'],
            marital_status=to_latn(d['marital_status']), marital_status_uz_latn=to_latn(d['marital_status']),
            marital_status_uz_cyrl=d['marital_status'], marital_status_ru=d['marital_status'],
            university=to_latn(d['university']), university_uz_latn=to_latn(d['university']),
            university_uz_cyrl=d['university'], university_ru=d['university'],
            specialty=to_latn(d['specialty']), specialty_uz_latn=to_latn(d['specialty']),
            specialty_uz_cyrl=d['specialty'], specialty_ru=d['specialty'],
            academic_degree=to_latn(d['academic_degree']), academic_degree_uz_latn=to_latn(d['academic_degree']),
            academic_degree_uz_cyrl=d['academic_degree'], academic_degree_ru=d['academic_degree'],
            training=to_latn(d['training']), training_uz_latn=to_latn(d['training']),
            training_uz_cyrl=d['training'], training_ru=d['training'],
            position=to_latn(d['position']), position_uz_latn=to_latn(d['position']),
            position_uz_cyrl=d['position'], position_ru=d['position'],
            career_level=to_latn(d['career_level']), career_level_uz_latn=to_latn(d['career_level']),
            career_level_uz_cyrl=d['career_level'], career_level_ru=d['career_level'],
            total_experience=to_latn(d['total_experience']), total_experience_uz_latn=to_latn(d['total_experience']),
            total_experience_uz_cyrl=d['total_experience'], total_experience_ru=d['total_experience'],
            leadership_experience=to_latn(d['leadership_experience']),
            leadership_experience_uz_latn=to_latn(d['leadership_experience']),
            leadership_experience_uz_cyrl=d['leadership_experience'],
            leadership_experience_ru=d['leadership_experience'],
            leadership_positions=to_latn(d['leadership_positions']),
            leadership_positions_uz_latn=to_latn(d['leadership_positions']),
            leadership_positions_uz_cyrl=d['leadership_positions'],
            leadership_positions_ru=d['leadership_positions'],
            health=d['health'],
            last_medical_treatment=to_latn(d['last_medical_treatment']),
            last_medical_treatment_uz_latn=to_latn(d['last_medical_treatment']),
            last_medical_treatment_uz_cyrl=d['last_medical_treatment'],
            last_medical_treatment_ru=d['last_medical_treatment'],
            medical_checkup=to_latn(d['medical_checkup']),
            medical_checkup_uz_latn=to_latn(d['medical_checkup']),
            medical_checkup_uz_cyrl=d['medical_checkup'],
            medical_checkup_ru=d['medical_checkup'],
            health_problems=to_latn(d['health_problems']),
            health_problems_uz_latn=to_latn(d['health_problems']),
            health_problems_uz_cyrl=d['health_problems'],
            health_problems_ru=d['health_problems'],
            description=to_latn(d['description']), description_uz_latn=to_latn(d['description']),
            description_uz_cyrl=d['description'], description_ru=d['description'],
            state_events=to_latn(d['state_events']), state_events_uz_latn=to_latn(d['state_events']),
            state_events_uz_cyrl=d['state_events'], state_events_ru=d['state_events'],
            is_active=True,
        )

        for i, fm in enumerate(d['family']):
            FamilyMember.objects.create(
                representative=rep,
                relation=fm['relation'], order=i + 1,
                name=to_latn(fm['name']), name_uz_latn=to_latn(fm['name']),
                name_uz_cyrl=fm['name'], name_ru=fm['name'],
                info=to_latn(fm['info']), info_uz_latn=to_latn(fm['info']),
                info_uz_cyrl=fm['info'], info_ru=fm['info'],
                note=to_latn(fm['note']), note_uz_latn=to_latn(fm['note']),
                note_uz_cyrl=fm['note'], note_ru=fm['note'],
            )

        return rep
